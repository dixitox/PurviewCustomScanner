import requests
import json
import re
import tiktoken
from azure.ai.inference.models import SystemMessage, UserMessage
from pyapacheatlas.core import AtlasClassification, AtlasEntity
from datetime import datetime
import os
import random
from pathlib import Path
import uuid
import PyPDF2
from pptx import Presentation
from docx import Document as DocxDocument
from custom_libs.custom_logging import get_logger

# OCR imports - with graceful fallback if not installed
try:
    from pdf2image import convert_from_path
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger = get_logger()
    logger.warning("OCR libraries not available. Install 'pytesseract', 'pdf2image', and 'Pillow' for scanned PDF support.")

logger = get_logger()

def filesystemFileSampleList(input_file_list, file_sample_size,file_system_path,):
    """
    From a list of files identified, generates a subset of files based on the given sample
    size, extracts the contents from each file, and returns a list.

    Args:
        input_file_list (list): A list of file paths (as strings).
        file_sample_size (int): The desired sample size of files to process.
        file_system_path (str): The file system path to be used as a parent object.

    Returns:
        list: A list of dictionaries containing file information. Each dictionary includes:
            - 'id': A unique identifier (UUID) for the file.
            - 'name': The name of the file.
            - 'created_datetime': The creation timestamp of the file.
            - 'created_by': The user ID of the file creator.
            - 'size': The file size in bytes.
            - 'last_modified_datetime': The last modification timestamp of the file.
            - 'last_modified_by': The user ID of the last modifier.
            - 'source': The original file path.

    Example:
        input_files = ['/path/to/file1.txt', '/path/to/file2.txt', '/path/to/file3.txt']
        sample_size = 2
        result = filesystemFileSampleList(input_files, sample_size)
        # Example output (for illustration purposes):
        # [{'id': '12345', 'name': 'file1.txt', 'created_datetime': ..., 'created_by': 'user1', ...},
        #  {'id': '67890', 'name': 'file3.txt', 'created_datetime': ..., 'created_by': 'user2', ...}]
    """
    num_files_found = len(input_file_list)
    files_to_process = []
    content = ''

    def build_file_info(path_str):
        fileFullPath = Path(path_str)
        fileName = fileFullPath.name
        root_name = Path(file_system_path).name
        try:
            rel_folder = fileFullPath.parent.relative_to(file_system_path).as_posix()
        except ValueError:
            rel_folder = fileFullPath.parent.as_posix()

        print(f"Processing {fileName}")
        local_content = ''
        if fileName.endswith(".docx"):
            try:
                document = DocxDocument(fileFullPath)
                local_content = "\n".join([paragraph.text for paragraph in document.paragraphs])
            except Exception as err:
                logger.error(f"Error processing document: {err}")
        elif fileName.endswith(".pdf"):
            try:
                local_content = extractPDFContent(fileFullPath)
            except Exception as err:
                logger.error(f"Error processing document: {err}")
        elif fileName.endswith(".pptx"):
            try:
                local_content = extractPPTXContent(fileFullPath)
            except Exception as err:
                logger.error(f"Error processing document: {err}")

        file_info = {
            'id': str(uuid.uuid4()),
            'name': fileName,
            'created_datetime': datetime.fromtimestamp(fileFullPath.stat().st_ctime),
            'created_by': str(fileFullPath.stat().st_uid),
            'size': fileFullPath.stat().st_size,
            'last_modified_datetime': datetime.fromtimestamp(fileFullPath.stat().st_mtime),
            'last_modified_by': str(fileFullPath.stat().st_uid),
            'source': path_str,
            'parentObject': f"{root_name}/{rel_folder}" if rel_folder else root_name,
            'typedef': 'FileSystemFile',
            'fs_root': root_name,
            'fs_folder': rel_folder,
            'content': local_content,
        }
        return file_info

    if num_files_found > file_sample_size:
        for item in random.sample(range(num_files_found), file_sample_size):
            files_to_process.append(build_file_info(input_file_list[item]))
    else:
        for item in input_file_list:
            files_to_process.append(build_file_info(item))

    return files_to_process

def sharepointFileSampleList(input_file_list, file_sample_size):
    """
    Generates a list of files to process from an input list of files.

    Args:
        input_file_list (list): A list of dictionaries containing file information.
            Each dictionary should have a 'name' key representing the file name.
        file_sample_size (int): The desired sample size of files to process.

    Returns:
        list: A list of file names to process. If the number of files found is greater
        than the sample size, a random sample of files is returned. Otherwise, all files
        in the input list are included.

    Example:
        input_files = [{'name': 'file1.txt'}, {'name': 'file2.txt'}, {'name': 'file3.txt'}]
        sample_size = 2
        result = sharepointFileSampleList(input_files, sample_size)
        # Example output: ['file1.txt', 'file3.txt']
    """
    num_files_found = len(input_file_list)
    files_to_process = []

    if num_files_found > file_sample_size:
        # Randomly sample files if more files are found than the desired sample size
        for ind in random.sample(range(num_files_found), file_sample_size):
            files_to_process.append(input_file_list[ind]['name'])
    else:
        # Include all files if the number of files found is less than or equal to the sample size
        for ind in input_file_list:
            files_to_process.append(ind['name'])

    return files_to_process

def listFilesystemFiles(mypath, extensions):
    """
    Lists files in a directory matching specified extensions.

    Args:
        mypath (str): Path to the directory. If empty, defaults to "SampleFiles".
        extensions (list): List of file extensions (e.g., ['txt', 'csv', 'jpg'] or ['.txt', '.csv', '.jpg']).

    Returns:
        list: List of filenames matching the specified extensions.
    """
    # Default to SampleFiles if path is empty
    if not mypath:
        mypath = "SampleFiles"
    
    # Ensure extensions have dots
    normalized_extensions = [ext if ext.startswith('.') else f'.{ext}' for ext in extensions]
    
    logger.info(f"Scanning filesystem: '{mypath}' for extensions: {normalized_extensions}")
    
    files_found = [
        os.path.join(dirpath, f) 
        for (dirpath, _, filenames) in os.walk(mypath) 
        for f in filenames 
        if any(f.endswith(ext) for ext in normalized_extensions)
    ]
    
    logger.info(f"Found {len(files_found)} files in filesystem")
    return files_found

def extractPDFContent(fileFullPath):
    # Function to convert PDF files to Text with enhanced extraction
    try:
        pdffileobj = open(fileFullPath,'rb')
        pdfreader = PyPDF2.PdfReader(pdffileobj)
        num_pages = len(pdfreader.pages)
        content = ''
        for page in range(num_pages):
            pageobj = pdfreader.pages[page]
            page_text = pageobj.extract_text()
            if page_text:
                content += page_text + '\n'
        
        # Clean up the extracted text
        content = content.strip()
        
        # If minimal content extracted, try OCR
        if len(content) < 50:
            logger.warning(f"Minimal text extracted from PDF ({len(content)} chars). Attempting OCR...")
            if OCR_AVAILABLE:
                ocr_content = extractPDFContentWithOCR(fileFullPath)
                if ocr_content and len(ocr_content) > len(content):
                    logger.info(f"OCR extracted {len(ocr_content)} chars (vs {len(content)} from standard extraction)")
                    return ocr_content
            else:
                logger.warning("OCR not available. Install pytesseract, pdf2image, and Pillow for scanned PDF support.")
            
        return content if content else None
    except Exception as err:
        logger.error(f"Error processing document: {err}")
        return None

def extractPDFContentWithOCR(fileFullPath, max_pages=10):
    """
    Extract text from PDF using OCR (Optical Character Recognition).
    Useful for scanned/image-based PDFs where standard text extraction fails.
    
    Args:
        fileFullPath: Path to the PDF file
        max_pages: Maximum number of pages to process (to avoid long processing times)
    
    Returns:
        Extracted text content or None if OCR fails
    """
    if not OCR_AVAILABLE:
        return None
        
    try:
        # Convert PDF pages to images
        logger.info(f"Converting PDF to images for OCR processing...")
        images = convert_from_path(fileFullPath, dpi=200, first_page=1, last_page=max_pages)
        
        content = ''
        for i, image in enumerate(images, 1):
            logger.info(f"  Processing page {i}/{len(images)} with OCR...")
            # Extract text from image using Tesseract OCR
            page_text = pytesseract.image_to_string(image, lang='eng')
            if page_text:
                content += page_text + '\n'
        
        content = content.strip()
        logger.info(f"OCR extraction completed: {len(content)} characters extracted from {len(images)} pages")
        return content if content else None
        
    except Exception as err:
        logger.error(f"OCR extraction failed: {err}")
        return None

def extractPPTXContent(fileFullPath):
    # Function to convert PPTX files to Text
    try:
        presentation = Presentation(fileFullPath)
        content = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + '\n'
        return content
    except Exception as err:
        logger.error(f"Error processing document: {err}")
        return None

def getAADToken(tenant_id: str, client_id: str, client_secret: str, resource_url: str):
    """
    Authenticates Service Principal to the provided Resource URL, and returns the OAuth Access Token
    """
    url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/token"
    payload= f'grant_type=client_credentials&client_id={client_id}&client_secret={client_secret}&resource={resource_url}'
    headers = {
    'Content-Type': 'application/x-www-form-urlencoded'
    }
    response = requests.request("POST", url, headers=headers, data=payload)
    access_token = json.loads(response.text)['access_token']
    return access_token

def moveCollection(collectionId: str, endpoint: str, token: str, guids):
    url = f"{endpoint}/datamap/api/entity/moveTo?api-version=2023-09-01&collectionId={collectionId}"
    payload={
        "entityGuids": guids
    }
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    payload_json = json.dumps(payload)
    response = requests.request("POST", url, headers=headers, data=payload_json)
    return json.loads(response.text)

def listAssets(endpoint: str, token: str, query: str):
    url = f"{endpoint}/datamap/api/search/query?api-version=2023-09-01"
    payload={
        "keywords": query
    }
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    payload_json = json.dumps(payload)
    response = requests.request("POST", url, headers=headers, data=payload_json)
    return json.loads(response.text)

def deleteAssets(purviewClient, endpoint: str, token: str, query: str):
    url = f"{endpoint}/datamap/api/search/query?api-version=2023-09-01"
    payload={
        "keywords": query
    }
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    payload_json = json.dumps(payload)
    response = requests.request("POST", url, headers=headers, data=payload_json)
    assets = [asset['id'] for asset in json.loads(response.text)['value']]
    delResponse = purviewClient.delete_entity(guid=assets)
    return json.dumps(delResponse, indent=2)

def estimateTokens(fileContents,textLength,classificationsStr,aiModel):
    """
    Estimates the total number of tokens in a given user input, including system messages and user prompts.

    Args:
        fileContents (List[Doc]): List of document contents to be analyzed.
        textLenght (int): Number of characters to be analyzed from of each document.
        classificationsStr (str): String containing the list of classifications to be evaluated against each document.

    Returns:
        totalTokens: Int containing the estimated total number of tokens.
    """
    totalTokens = 0
    for ind in range(len(fileContents)):
        fileContent = fileContents[ind]
        # Text extracted from scanned documents up to the number of characters specified by textLength
        # Handle None content gracefully
        text = (fileContent.get('content') or '')[:textLength]
        
        # Count meaningful characters
        if re.search(r'[a-zA-Z0-9]', text):
            # Generate user prompt
            userPrompt = (
                "role: system, content: You are an AI assistant that helps people classify information contained in documents."
                +"role: user, content: "
                +"Classify the following text: "+text
                +"Using only any of the following classifications: \n"+classificationsStr
                +"The answer should only return one of the listed classifications."
            )
            # enc = tiktoken.get_encoding("cl100k_base")
            # enc = tiktoken.encoding_for_model('gpt-4')
            enc = tiktoken.encoding_for_model(aiModel)
            numTokens = len(enc.encode(userPrompt))
        else:
            numTokens = 0
        totalTokens += numTokens
    return totalTokens

def unstructuredDataClassification(fileContents,textLength,llmClient,llmodel,classificationsStr):
    """
    This function will run the contents of each file contained in the subset against the 
    Azure Open AI large language model (GPT-4). The result will be a list of dictionaries
    containing file metadata and also the classification assigned by the LLM model.
    Enhanced with filename-based fallback and better content detection.
    """
    for ind in range(len(fileContents)):
        fileContent = fileContents[ind]
        filename = fileContent.get('name', '')
        # Text extracted from documents up to the number of characters specified by textLength
        text = (fileContent.get('content') or '')[:textLength]
        
        # Count meaningful content (alphanumeric characters)
        meaningful_chars = len(re.findall(r'[a-zA-Z0-9]', text))
        
        # If we have sufficient meaningful content (at least 20 chars)
        if meaningful_chars >= 20:
            # Generate enhanced user prompt
            userPrompt = (
                f"Classify this document excerpt. Filename: {filename}\n\n"
                f"Content:\n{text}\n\n"
                f"Choose ONE classification from this list:\n{classificationsStr}\n"
                "Return ONLY the classification name, nothing else."
            )
            # Submit prompts to LLM model with lower temperature for more consistent results
            response = llmClient.complete(
                messages=[
                    SystemMessage(content="You are an insurance specialist that helps people classify information contained in documents. Use both the filename and content to determine the most accurate classification."),
                    UserMessage(content=userPrompt),
                ],
                max_tokens=4096,
                temperature=0.3,  # Lower temperature for more deterministic classification
                top_p=1.0,
                model=llmodel
            )
            classification = response.choices[0].message.content.strip()
            print(f"Processing {filename} -> {classification} (content: {meaningful_chars} chars)")
            fileContents[ind].update({'classification':classification})
            # Comment out the line below to determine the exact number of tokens consumed by each document
            print(response.usage)
        else:
            # Fallback: Try to infer from filename
            classification = classify_from_filename(filename, classificationsStr)
            print(f"Processing {filename} -> {classification} (insufficient content, classified by filename)")
            fileContents[ind].update({'classification':classification})
    return fileContents

def classify_from_filename(filename, classificationsStr):
    """
    Attempt to classify a document based on its filename when content extraction fails.
    Returns the most likely classification or 'Other' if no match.
    """
    filename_lower = filename.lower()
    classifications = [c.strip() for c in classificationsStr.strip().split('\n') if c.strip()]
    
    # Keyword patterns for each classification type
    patterns = {
        'Insurance Claim': ['claim', 'loss', 'damage', 'incident'],
        'Insurance Policy': ['policy', 'insurance', 'coverage', 'wording'],
        'Report': ['report', 'inspection', 'assessment', 'analysis', 'toxicology', 'mold', 'asbestos'],
        'Invoice': ['invoice', 'bill', 'receipt', 'payment'],
        'Sales Receipt': ['receipt', 'sales', 'purchase'],
        'PII': ['personal', 'pii', 'confidential', 'ssn'],
    }
    
    # Check filename against patterns
    for classification, keywords in patterns.items():
        if classification in classifications:
            if any(keyword in filename_lower for keyword in keywords):
                return classification
    
    # Default to 'Other' if it exists, otherwise 'Empty Content'
    if 'Other' in classifications:
        return 'Other'
    elif 'Empty Content' in classifications:
        return 'Empty Content'
    else:
        return classifications[0] if classifications else 'Unknown'

def loadPurviewAssets(purviewClient,allFileContent):
    """
    Load assets for SharePoint and FileSystem data sources creating hierarchical structure:
    - SharePoint: Account -> RootFolder -> Folder(s) -> Files  
    - FileSystem: Root -> Folder(s) -> Files
    Handles recursive folder structures, creating assets for all discovered folders.
    """
    if not allFileContent:
        return {"all": [], "file": []}

    batchEntities = []
    tempFileGuids = []  # Track file temp GUIDs to map after upload
    newGuid = -1000
    folderEntitiesCreated = {}  # Track created folders to avoid duplicates

    sample = allFileContent[0]
    typedef = sample.get("typedef")
    
    print(f"🔍 DEBUG: typedef detected = '{typedef}'")
    print(f"🔍 DEBUG: sample keys = {list(sample.keys())}")
    print(f"🔍 DEBUG: Total files to process = {len(allFileContent)}")

    if typedef in ["SharepointFile", "SharePoint"]:
        print(f"✅ Creating SharePoint hierarchy with recursive folder support")
        
        # Collect unique folders from all files
        account = sample.get("sharepoint_account") or (sample.get("source", "").split("/")[2] if sample.get("source") else "")
        root_folder = sample.get("sharepoint_root", "")
        
        # Create Account entity (once)
        accountGuid = newGuid
        accountEntity = AtlasEntity(
            name=account,
            typeName="SharepointAccount",
            qualified_name=f"customScanner://sharepoint/{account}",
            guid=accountGuid,
            attributes={
                "name": account,
                "qualifiedName": f"customScanner://sharepoint/{account}",
            }
        )
        batchEntities.append(accountEntity)
        newGuid -= 1
        print(f"  📁 Account: {account}")

        # Create Root Folder entity (once) with relationship to Account
        rootGuid = newGuid
        rootEntity = AtlasEntity(
            name=root_folder,
            typeName="SharepointRootFolder",
            qualified_name=f"customScanner://sharepoint/{account}/{root_folder}",
            guid=rootGuid,
            attributes={
                "name": root_folder,
                "qualifiedName": f"customScanner://sharepoint/{account}/{root_folder}",
                "account": {
                    "guid": accountGuid,
                    "typeName": "SharepointAccount",
                    "qualifiedName": f"customScanner://sharepoint/{account}"
                }
            },
            relationshipAttributes={
                "account": {
                    "guid": accountGuid,
                    "typeName": "SharepointAccount",
                    "qualifiedName": f"customScanner://sharepoint/{account}"
                }
            }
        )
        batchEntities.append(rootEntity)
        newGuid -= 1
        print(f"  📁 Root Folder: {root_folder}")

        # Collect all unique folders from files
        unique_folders = set()
        for file in allFileContent:
            folder_path = (file.get("sharepoint_folder") or "").strip("/")
            if folder_path:
                unique_folders.add(folder_path)
        
        # Create folder entities for each unique folder with parent relationships
        for folder_path in sorted(unique_folders):
            if folder_path not in folderEntitiesCreated:
                folder_name = folder_path.split("/")[-1] if "/" in folder_path else folder_path
                folderGuid = newGuid
                
                # Determine parent (either root or parent folder)
                parent_relationship = {}
                if "/" in folder_path:
                    # Has parent folder
                    parent_path = "/".join(folder_path.split("/")[:-1])
                    if parent_path in folderEntitiesCreated:
                        parent_relationship = {
                            "parentFolder": {
                                "guid": folderEntitiesCreated[parent_path],
                                "typeName": "SharepointFolder",
                                "qualifiedName": f"customScanner://sharepoint/{account}/{root_folder}/{parent_path}"
                            }
                        }
                    else:
                        # Parent folder not yet created, link to root
                        parent_relationship = {
                            "rootFolder": {
                                "guid": rootGuid,
                                "typeName": "SharepointRootFolder",
                                "qualifiedName": f"customScanner://sharepoint/{account}/{root_folder}"
                            }
                        }
                else:
                    # Top-level folder, parent is root
                    parent_relationship = {
                        "rootFolder": {
                            "guid": rootGuid,
                            "typeName": "SharepointRootFolder",
                            "qualifiedName": f"customScanner://sharepoint/{account}/{root_folder}"
                        }
                    }
                
                # Add parent to both attributes and relationshipAttributes
                folder_attributes = {
                    "name": folder_name,
                    "qualifiedName": f"customScanner://sharepoint/{account}/{root_folder}/{folder_path}",
                    "path": folder_path,
                }
                # Add parent reference to attributes for hierarchy
                if parent_relationship:
                    folder_attributes.update(parent_relationship)
                
                folderEntity = AtlasEntity(
                    name=folder_name,
                    typeName="SharepointFolder",
                    qualified_name=f"customScanner://sharepoint/{account}/{root_folder}/{folder_path}",
                    guid=folderGuid,
                    attributes=folder_attributes,
                    relationshipAttributes=parent_relationship
                )
                batchEntities.append(folderEntity)
                folderEntitiesCreated[folder_path] = folderGuid
                newGuid -= 1
                print(f"  📁 Folder: {folder_path}")
        
        print(f"  📄 Total files: {len(allFileContent)}")

        # Files with parent folder relationships
        for file in allFileContent:
            fileGuid = newGuid
            file_qualified = file.get("source") or f"customScanner://sharepoint/{account}/{root_folder}/{folder_path}/{file['name']}"
            
            # Convert datetime objects to ISO format strings
            created_date = file.get("created_datetime")
            if hasattr(created_date, 'isoformat'):
                created_date = created_date.isoformat()
            modified_date = file.get("last_modified_datetime")
            if hasattr(modified_date, 'isoformat'):
                modified_date = modified_date.isoformat()
            
            # Determine parent folder for this file
            folder_path = (file.get("sharepoint_folder") or "").strip("/")
            file_relationship = {}
            if folder_path and folder_path in folderEntitiesCreated:
                # File belongs to a folder
                file_relationship = {
                    "folder": {
                        "guid": folderEntitiesCreated[folder_path],
                        "typeName": "SharepointFolder",
                        "qualifiedName": f"customScanner://sharepoint/{account}/{root_folder}/{folder_path}"
                    }
                }
            else:
                # File is in root folder
                file_relationship = {
                    "rootFolder": {
                        "guid": rootGuid,
                        "typeName": "SharepointRootFolder",
                        "qualifiedName": f"customScanner://sharepoint/{account}/{root_folder}"
                    }
                }
            
            # Add parent to both attributes and relationshipAttributes
            file_attributes = {
                "size": file.get("size"),
                "createdBy": file.get("created_by"),
                "createdDate": created_date,
                "lastModifiedDate": modified_date,
                "lastModifiedBy": file.get("last_modified_by"),
            }
            # Add parent reference to attributes for hierarchy
            if file_relationship:
                file_attributes.update(file_relationship)
            
            fileEntity = AtlasEntity(
                name=file["name"],
                typeName="SharepointFile",
                qualified_name=file_qualified,
                guid=fileGuid,
                attributes=file_attributes,
                relationshipAttributes=file_relationship
            )
            batchEntities.append(fileEntity)
            tempFileGuids.append(fileGuid)
            newGuid -= 1

    elif typedef == "FileSystemFile":
        print(f"✅ Creating FileSystem hierarchy with recursive folder support")
        
        # Get root from first file
        root = sample.get("fs_root", "FileSystem")
        
        # Create Root entity (once)
        rootGuid = newGuid
        rootEntity = AtlasEntity(
            name=root,
            typeName="FileSystemRoot",
            qualified_name=f"customScanner://filesystem/{root}",
            guid=rootGuid,
            attributes={
                "name": root,
                "qualifiedName": f"customScanner://filesystem/{root}",
            }
        )
        batchEntities.append(rootEntity)
        newGuid -= 1
        print(f"  📁 Root: {root}")

        # Collect all unique folders from files
        unique_folders = set()
        for file in allFileContent:
            folder_path = (file.get("fs_folder") or "").strip("/")
            if folder_path and folder_path != ".":
                unique_folders.add(folder_path)
        
        # Create folder entities for each unique folder with parent relationships
        for folder_path in sorted(unique_folders):
            if folder_path not in folderEntitiesCreated:
                folder_name = folder_path.split("/")[-1] if "/" in folder_path else folder_path
                folderGuid = newGuid
                
                # Determine parent (either root or parent folder)
                parent_relationship = {}
                if "/" in folder_path:
                    # Has parent folder
                    parent_path = "/".join(folder_path.split("/")[:-1])
                    if parent_path in folderEntitiesCreated:
                        parent_relationship = {
                            "parentFolder": {
                                "guid": folderEntitiesCreated[parent_path],
                                "typeName": "FileSystemFolder",
                                "qualifiedName": f"customScanner://filesystem/{root}/{parent_path}"
                            }
                        }
                    else:
                        # Parent folder not yet created, link to root
                        parent_relationship = {
                            "root": {
                                "guid": rootGuid,
                                "typeName": "FileSystemRoot",
                                "qualifiedName": f"customScanner://filesystem/{root}"
                            }
                        }
                else:
                    # Top-level folder, parent is root
                    parent_relationship = {
                        "root": {
                            "guid": rootGuid,
                            "typeName": "FileSystemRoot",
                            "qualifiedName": f"customScanner://filesystem/{root}"
                        }
                    }
                
                # Add parent to both attributes and relationshipAttributes
                folder_attributes = {
                    "name": folder_name,
                    "qualifiedName": f"customScanner://filesystem/{root}/{folder_path}",
                    "path": folder_path,
                }
                # Add parent reference to attributes for hierarchy
                if parent_relationship:
                    folder_attributes.update(parent_relationship)
                
                folderEntity = AtlasEntity(
                    name=folder_name,
                    typeName="FileSystemFolder",
                    qualified_name=f"customScanner://filesystem/{root}/{folder_path}",
                    guid=folderGuid,
                    attributes=folder_attributes,
                    relationshipAttributes=parent_relationship
                )
                batchEntities.append(folderEntity)
                folderEntitiesCreated[folder_path] = folderGuid
                newGuid -= 1
                print(f"  📁 Folder: {folder_path}")
        
        print(f"  📄 Total files: {len(allFileContent)}")

        # Files with parent folder relationships
        for file in allFileContent:
            fileGuid = newGuid
            file_qualified = file.get("source") or f"customScanner://filesystem/{root}/{folder_path}/{file['name']}"
            
            # Convert datetime objects to ISO format strings
            created_date = file.get("created_datetime")
            if hasattr(created_date, 'isoformat'):
                created_date = created_date.isoformat()
            modified_date = file.get("last_modified_datetime")
            if hasattr(modified_date, 'isoformat'):
                modified_date = modified_date.isoformat()
            
            # Determine parent folder for this file
            folder_path = (file.get("fs_folder") or "").strip("/")
            file_relationship = {}
            if folder_path and folder_path != "." and folder_path in folderEntitiesCreated:
                # File belongs to a folder
                file_relationship = {
                    "folder": {
                        "guid": folderEntitiesCreated[folder_path],
                        "typeName": "FileSystemFolder",
                        "qualifiedName": f"customScanner://filesystem/{root}/{folder_path}"
                    }
                }
            else:
                # File is in root folder
                file_relationship = {
                    "root": {
                        "guid": rootGuid,
                        "typeName": "FileSystemRoot",
                        "qualifiedName": f"customScanner://filesystem/{root}"
                    }
                }
            
            # Add parent to both attributes and relationshipAttributes
            file_attributes = {
                "size": file.get("size"),
                "createdBy": file.get("created_by"),
                "createdDate": created_date,
                "lastModifiedDate": modified_date,
                "lastModifiedBy": file.get("last_modified_by"),
            }
            # Add parent reference to attributes for hierarchy
            if file_relationship:
                file_attributes.update(file_relationship)
            
            fileEntity = AtlasEntity(
                name=file["name"],
                typeName="FileSystemFile",
                qualified_name=file_qualified,
                guid=fileGuid,
                attributes=file_attributes,
                relationshipAttributes=file_relationship
            )
            batchEntities.append(fileEntity)
            tempFileGuids.append(fileGuid)
            newGuid -= 1
    else:
        # Fallback for any other dataset types
        print(f"⚠️  Using fallback (single parent entity) - typedef: {typedef}")
        file = sample
        parentGuid = newGuid
        newEntity = AtlasEntity(
            name=file["parentObject"],
            typeName=file["typedef"],
            qualified_name=f"customScanner://{file['parentObject']}",
            guid=parentGuid,
        )
        batchEntities.append(newEntity)
        tempFileGuids.append(parentGuid)

    # Convert the individual entities into json before uploading.
    results = purviewClient.upload_entities(
        batch=batchEntities,
        batch_size=20
    )

    guidAssignments = results.get("guidAssignments", {})
    allGuids = [v for v in guidAssignments.values()]
    fileGuids = [guidAssignments[str(guid)] for guid in tempFileGuids if str(guid) in guidAssignments]

    return {"all": allGuids, "file": fileGuids}

def rollupClassifications(allFileContent):
    """
    Obtain parent level classifications to represent the classifications of the individual
    files stored in the data source.
    """
    classificationsFound = list(set([file['classification'] for file in allFileContent]))
    return classificationsFound

def applyPurviewClassifications(purviewClient,guids,classificationsFound):
    """
    Apply classifications to all provided GUIDs
    """
    atlasClassifications = [AtlasClassification(classification_name) for classification_name in classificationsFound]
    results = []
    for guid in guids:
        resp = purviewClient.classify_entity(
            guid=guid,
            classifications=atlasClassifications,
            force_update=True
        )
        results.append(resp)
    return results